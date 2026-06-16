#!/usr/bin/env python3
"""Generate demo-support collateral for Vela OS, into public/downloads/ so the Worker
serves them (and the /guide page can link them as downloads):
  - public/downloads/Vela-OS-Demo-Deck.pptx   (dark, branded slide deck)
  - public/downloads/Vela-OS-One-Pager.pdf    (clean light executive brief)
Content mirrors public/guide.html and the real proof data in public/data.jsx.
  python3 scripts/build-explainers.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from fpdf import FPDF
from fpdf.enums import XPos, YPos
MC = dict(new_x=XPos.LMARGIN, new_y=YPos.NEXT)  # keep cursor at left margin after multi_cell

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "..", "public", "downloads")
os.makedirs(OUT, exist_ok=True)

LIVE = "https://roai.emergedigital.ae"
GUIDE = LIVE + "/guide"

# ---------- shared content ----------
TAGLINE = "The operating system for AI-first customer experience."
SUB = "Vela OS turns FPT's full CX offering into one measurable surface - a dual-mode agency command center and customer portal, where every AI engagement is priced, delivered, and proven in ROAI."
ROAI_F = "ROAI = (Value delivered - AI cost) / AI cost"
ROAI_NOTE = "A 5.0x ROAI means every $1 of AI spend returned $5 of measured value - hours saved plus revenue uplift."
STATS = [("2-3x", "higher revenue growth for CX leaders"),
         ("~70%", "rate experience as important as product or price"),
         ("~60%", "of interactions span multiple channels")]
PILLARS = [
    ("Experience Strategy", "Maturity assessment, marketecture, journey analysis"),
    ("Experience Design", "Research, adaptive UI, design systems, usability"),
    ("Experience Platform", "DXP build/re-platform, MarTech, DX engineering"),
    ("Commerce Platform", "Composable commerce, B2B/B2C, accelerators"),
    ("Experience Insights", "CDP & Customer 360, analytics portals, AI agents"),
    ("Run & Optimize", "L1-L3 maintenance, 24x7 monitoring, ops support"),
]
AGENTS = ["ON.Optima - Answer Engine Optimization", "ON.X - Agentic Assistant",
          "ON.E - Commerce Accelerator", "ON.Match - AI Recommendations", "ON.Browser - Agentic Browser"]
CASES = [
    ("ON.Optima - AEO", "+200% click-through rate", "Global technology leader",
     "Site ranking #6 (+18) - CTR 3.2% - featured snippets 30%"),
    ("ON.Optima - AEO", "Cited across 120+ markets", "Global health F&B leader",
     "120+ markets - consistent AI answer boxes - acquisition cost down"),
    ("ON.E - Commerce", "MVP in 2 months, not 10", "KSA megaproject (PIF)",
     "Time to MVP 2 mo (vs 10) - build time -55% - requirements -40%"),
    ("Data & AI CoE", "Campaign cycles: weeks to minutes", "FPT x Salesforce ASEAN CoE",
     "First in ASEAN - weeks to minutes - ~100 Agentforce certs by 2026"),
]
PARTNERS = ["Adobe (Specialized)", "Sitecore (Platinum APAC)", "Salesforce (First ASEAN Data & AI CoE)",
            "Liferay (Global Platinum)", "Shopify (Plus)"]
DEMO = [
    ("Agency overview", "Blended ROAI, revenue and margin across all clients.", LIVE + "/?view=agency&section=home"),
    ("ROAI analytics", "The formula, the trend, the per-pillar breakdown.", LIVE + "/?view=agency&section=analytics"),
    ("Client deep dive", "Drill into one account + its real case result.", LIVE + "/?view=agency&client=helios"),
    ("Billing -> Live", "Real Metronome invoices and commit drawdown.", LIVE + "/?view=agency&section=billing&tab=live"),
    ("Customer portal", "The same engagement, the client's own view.", LIVE + "/?view=client&client=helios&section=roai"),
    ("AI marketplace", "Launch a new ON.Ecosystem agent in a click.", LIVE + "/?view=client&client=helios&section=marketplace"),
    ("Customer billing & usage", "Credits, usage drawdown and top-up.", LIVE + "/?view=client&client=helios&section=billing&tab=usage"),
]

# ---------- palette ----------
BG = RGBColor(0x0C, 0x0C, 0x13)
PANEL = RGBColor(0x15, 0x15, 0x1E)
PANEL2 = RGBColor(0x1B, 0x1B, 0x26)
WHITE = RGBColor(0xF3, 0xF3, 0xF6)
MUTE = RGBColor(0x9C, 0x9C, 0xAA)
FAINT = RGBColor(0x6B, 0x6B, 0x77)
INDIGO = RGBColor(0x6366F1 >> 16 & 255, 0x6366F1 >> 8 & 255, 0x6366F1 & 255)
INDIGO_FG = RGBColor(0xA5, 0xB4, 0xFC)
EMER = RGBColor(0x34, 0xD3, 0x99)
FONT = "Arial"

# ============================================================ PPTX ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h, fill=None, line=None, radius=True):
    shp = s.shapes.add_shape(5 if radius else 1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = space
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = FONT
    return tb


def eyebrow(s, t, x=0.9, y=0.62):
    txt(s, x, y, 11, 0.4, [[(t.upper(), 12, INDIGO_FG, True)]])


def heading(s, t, x=0.9, y=1.0, size=33, w=11.5):
    txt(s, x, y, w, 1.2, [[(t, size, WHITE, True)]])


def logo_mark(s, x, y, sz=0.42):
    sq = box(s, x, y, sz, sz, fill=INDIGO)
    txt(s, x, y - 0.02, sz, sz, [[("V", int(sz * 40), WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def footer(s):
    txt(s, 0.9, 7.02, 9, 0.3, [[("Vela OS  ", 9, MUTE, True), ("Emerge Digital x FPT CX Services", 9, FAINT, False)]])
    txt(s, 9.5, 7.02, 2.93, 0.3, [[(LIVE.replace("https://", ""), 9, FAINT, False)]], align=PP_ALIGN.RIGHT)


# --- Slide 1: title ---
s = slide()
box(s, 0, 0, 13.333, 7.5, fill=None)
logo_mark(s, 0.9, 0.85, 0.6)
txt(s, 1.65, 0.86, 8, 0.7, [[("Vela ", 22, WHITE, True), ("OS", 22, INDIGO_FG, True)]])
txt(s, 1.65, 1.32, 8, 0.4, [[("FPT CX Services  -  ROAI Analytics", 12, MUTE, False)]])
txt(s, 0.9, 2.7, 11.5, 2.2, [[("The operating system for", 46, WHITE, True)], [("AI-first customer experience.", 46, WHITE, True)]], space=1.02)
txt(s, 0.9, 4.95, 10.8, 1.2, [[(SUB, 16, MUTE, False)]], space=1.12)
box(s, 0.9, 6.25, 6.4, 0.62, fill=PANEL, line=RGBColor(0x33, 0x33, 0x44))
txt(s, 1.12, 6.36, 6, 0.45, [[(ROAI_F, 13, WHITE, True)]])
txt(s, 8.0, 6.3, 4.4, 0.6, [[("Live, interactive demo", 13, EMER, True)], [("agency + customer portal", 11, MUTE, False)]], align=PP_ALIGN.RIGHT)

# --- Slide 2: opportunity ---
s = slide()
eyebrow(s, "The opportunity")
heading(s, "CX leaders win on growth, not just experience.")
txt(s, 0.9, 1.75, 11.4, 0.6, [[("Buyers expect AI-grade, cross-channel experiences. The gap is proof - most AI spend can't show what it returned. Vela OS closes that gap.", 14, MUTE, False)]], space=1.1)
cw, gap = 3.7, 0.35
for i, (n, l) in enumerate(STATS):
    x = 0.9 + i * (cw + gap)
    box(s, x, 2.7, cw, 2.0, fill=PANEL, line=RGBColor(0x2A, 0x2A, 0x36))
    txt(s, x + 0.3, 3.0, cw - 0.6, 0.9, [[(n, 44, INDIGO_FG, True)]])
    txt(s, x + 0.3, 3.95, cw - 0.6, 0.7, [[(l, 14, MUTE, False)]], space=1.05)
footer(s)

# --- Slide 3: two modes ---
s = slide()
eyebrow(s, "What it is")
heading(s, "One platform, two modes.")
modes = [("AGENCY COMMAND CENTER", INDIGO_FG, "Run the whole portfolio",
          ["Portfolio ROAI, revenue and margin", "Client directory with health tiers", "ROAI analytics per $ of AI invested", "Billing + live Metronome usage data"]),
         ("CUSTOMER PORTAL", EMER, "Show the customer their outcomes",
          ["A personal ROAI Center", "Active projects across six pillars", "AI marketplace to launch agents", "Self-serve billing, usage & top-ups"])]
for i, (tag, col, h, items) in enumerate(modes):
    x = 0.9 + i * 6.1
    box(s, x, 2.0, 5.7, 4.4, fill=PANEL, line=RGBColor(0x2A, 0x2A, 0x36))
    txt(s, x + 0.4, 2.35, 5, 0.3, [[(tag, 11, col, True)]])
    txt(s, x + 0.4, 2.75, 5, 0.5, [[(h, 20, WHITE, True)]])
    rows = [[("-  " + it, 14, MUTE, False)] for it in items]
    txt(s, x + 0.4, 3.5, 5, 2.7, rows, space=1.5)
footer(s)

# --- Slide 4: pillars ---
s = slide()
eyebrow(s, "The offering")
heading(s, "Six CX pillars, powered by the ON.Ecosystem.")
cw, ch, gx, gy = 3.7, 1.25, 0.35, 0.3
for i, (name, desc) in enumerate(PILLARS):
    r, c = divmod(i, 3)
    x = 0.9 + c * (cw + gx); y = 1.95 + r * (ch + gy)
    box(s, x, y, cw, ch, fill=PANEL, line=RGBColor(0x2A, 0x2A, 0x36))
    txt(s, x + 0.28, y + 0.2, cw - 0.5, 0.4, [[(name, 14.5, WHITE, True)]])
    txt(s, x + 0.28, y + 0.62, cw - 0.5, 0.6, [[(desc, 11.5, MUTE, False)]], space=1.05)
txt(s, 0.9, 5.05, 11.4, 0.3, [[("AGENTIC PRODUCTS", 11, INDIGO_FG, True)]])
txt(s, 0.9, 5.4, 11.5, 1.2, [[(a + ("        " if (i + 1) % 2 else "\n"), 13, MUTE, False) for i, a in enumerate(AGENTS)]], space=1.3)
footer(s)

# --- Slide 5: ROAI ---
s = slide()
eyebrow(s, "The differentiator")
heading(s, "Every engagement, measured in ROAI.")
box(s, 0.9, 2.0, 5.4, 3.8, fill=PANEL, line=RGBColor(0x2A, 0x2A, 0x36))
txt(s, 1.25, 2.35, 4.7, 0.3, [[("EXAMPLE CLIENT - BLENDED ROAI", 11, FAINT, True)]])
txt(s, 1.25, 2.75, 4.7, 1.3, [[("5.4x", 70, EMER, True)]])
txt(s, 1.25, 4.2, 4.7, 0.6, [[("$1 invested  ->  $5.40 of measured value returned", 13, MUTE, False)]], space=1.1)
txt(s, 6.7, 2.2, 5.7, 0.9, [[(ROAI_F, 17, WHITE, True)]])
txt(s, 6.7, 3.0, 5.7, 0.6, [[("Value = hours saved + revenue uplift, measured per engagement.", 13, MUTE, False)]], space=1.1)
txt(s, 6.7, 3.9, 5.7, 1.8, [[("FPT and the client always see the same source of truth - so renewals and expansion become a data conversation, not a debate.", 14, MUTE, False)]], space=1.15)
footer(s)

# --- Slide 6: live billing ---
s = slide()
eyebrow(s, "Not a mockup")
heading(s, "Live usage-based billing, wired to Metronome.")
txt(s, 0.9, 1.75, 11.4, 0.7, [[("The Billing -> Live panel reads a real Metronome account through a secure, read-only proxy - the API key never reaches the browser.", 14, MUTE, False)]], space=1.1)
cols = [("WHAT'S LIVE", INDIGO_FG, ["7 customers: prepaid-commit + hybrid models", "Finalized + draft invoices, 6 months history", "Commit/credit balances drawing down on usage", "Full catalog: metrics, products, rate cards, packages"]),
        ("WHY IT MATTERS", EMER, ["Usage-based pricing + committed-spend drawdown", "Government Net-60 / PO and self-serve top-ups", "Runs on production billing infrastructure", "The same proxy pattern drops into a real deployment"])]
for i, (tag, col, items) in enumerate(cols):
    x = 0.9 + i * 6.1
    box(s, x, 2.7, 5.7, 3.5, fill=PANEL, line=RGBColor(0x2A, 0x2A, 0x36))
    txt(s, x + 0.4, 3.0, 5, 0.3, [[(tag, 11, col, True)]])
    txt(s, x + 0.4, 3.45, 5, 2.6, [[("-  " + it, 13, MUTE, False)] for it in items], space=1.5)
footer(s)

# --- Slide 7: proof ---
s = slide()
eyebrow(s, "Proof")
heading(s, "Real FPT outcomes behind the demo.")
cw, ch, gx, gy = 5.7, 1.5, 0.4, 0.3
for i, (sol, head, who, metrics) in enumerate(CASES):
    r, c = divmod(i, 2)
    x = 0.9 + c * (cw + gx); y = 1.9 + r * (ch + gy)
    box(s, x, y, cw, ch, fill=PANEL, line=RGBColor(0x2A, 0x2A, 0x36))
    txt(s, x + 0.32, y + 0.16, cw - 0.6, 0.3, [[(sol, 10.5, INDIGO_FG, True)]])
    txt(s, x + 0.32, y + 0.44, cw - 0.6, 0.4, [[(head, 16, WHITE, True)]])
    txt(s, x + 0.32, y + 0.82, cw - 0.6, 0.3, [[(who, 10.5, FAINT, False)]])
    txt(s, x + 0.32, y + 1.08, cw - 0.6, 0.35, [[(metrics, 10, MUTE, False)]])
txt(s, 0.9, 5.45, 11.4, 0.3, [[("PARTNERS:  ", 11, INDIGO_FG, True), ("  ".join(PARTNERS), 11, MUTE, False)]], space=1.1)
txt(s, 0.9, 5.95, 11.4, 0.4, [[("1,000+ certified engineers      1,500+ platform certifications", 14, WHITE, True)]])
footer(s)

# --- Slide 8: demo flow ---
s = slide()
eyebrow(s, "Run the demo")
heading(s, "A 7-step walkthrough (~6-8 min).")
for i, (h, d, _u) in enumerate(DEMO):
    r, c = divmod(i, 2)
    x = 0.9 + c * 6.1; y = 1.95 + r * 1.02
    box(s, x, y, 0.55, 0.55, fill=INDIGO)
    txt(s, x, y - 0.02, 0.55, 0.55, [[(str(i + 1), 18, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 0.75, y - 0.02, 5.1, 0.4, [[(h, 14.5, WHITE, True)]])
    txt(s, x + 0.75, y + 0.34, 5.1, 0.5, [[(d, 11.5, MUTE, False)]], space=1.0)
txt(s, 0.9, 6.5, 11.4, 0.4, [[("Step-by-step with deep links + talk track:  ", 12, MUTE, False), (GUIDE, 12, INDIGO_FG, True)]])
footer(s)

# --- Slide 9: close ---
s = slide()
logo_mark(s, 0.9, 2.4, 0.7)
txt(s, 0.9, 3.35, 11.5, 1.0, [[("See it live.", 44, WHITE, True)]])
txt(s, 0.9, 4.35, 11.5, 0.5, [[("Interactive demo + customer portal, with live billing data.", 16, MUTE, False)]])
box(s, 0.9, 5.1, 6.0, 0.7, fill=PANEL, line=RGBColor(0x33, 0x33, 0x44))
txt(s, 1.15, 5.27, 5.6, 0.4, [[("Demo:  ", 13, FAINT, False), (LIVE.replace("https://", ""), 13, INDIGO_FG, True)]])
box(s, 0.9, 5.95, 6.0, 0.7, fill=PANEL, line=RGBColor(0x33, 0x33, 0x44))
txt(s, 1.15, 6.12, 5.6, 0.4, [[("Guide: ", 13, FAINT, False), (GUIDE.replace("https://", ""), 13, INDIGO_FG, True)]])
txt(s, 0.9, 7.02, 11.5, 0.3, [[("Built by Emerge Digital for FPT CX Services", 10, FAINT, False)]])

pptx_path = os.path.join(OUT, "Vela-OS-Demo-Deck.pptx")
prs.save(pptx_path)
print("wrote", pptx_path, f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")

# ============================================================ PDF ============================================================
IND = (79, 70, 229)
INK = (21, 21, 28)
GREY = (91, 91, 103)
LITE = (243, 243, 247)
LINE = (224, 224, 230)
EM = (5, 150, 105)


def sanitize(t):
    repl = {"–": "-", "—": "-", "−": "-", "→": "->", "•": "-",
            "≈": "~", "×": "x", "’": "'", "“": '"', "”": '"', "…": "..."}
    for k, v in repl.items():
        t = t.replace(k, v)
    return t


class PDF(FPDF):
    def header(self):
        self.set_fill_color(*IND); self.rect(0, 0, 210, 22, "F")
        self.set_xy(14, 6); self.set_text_color(255, 255, 255)
        self.set_font("Helvetica", "B", 15); self.cell(0, 5, "Vela OS", 0, 1)
        self.set_xy(14, 12); self.set_font("Helvetica", "", 9)
        self.cell(0, 5, "FPT CX Services  -  ROAI Analytics", 0, 1)
        self.set_xy(120, 8); self.set_font("Helvetica", "", 8.5)
        self.cell(76, 5, "Interactive demo + customer portal", 0, 1, "R")

    def footer(self):
        self.set_y(-13); self.set_text_color(*GREY); self.set_font("Helvetica", "", 8)
        self.cell(0, 5, sanitize("Built by Emerge Digital for FPT CX Services  -  " + LIVE.replace("https://", "")), 0, 0, "L")
        self.cell(0, 5, "Page %s" % self.page_no(), 0, 0, "R")


def h2(pdf, t):
    pdf.ln(2); pdf.set_text_color(*IND); pdf.set_font("Helvetica", "B", 13)
    pdf.cell(0, 7, sanitize(t), 0, 1); pdf.set_text_color(*INK)


def body(pdf, t, size=10):
    pdf.set_text_color(*GREY); pdf.set_font("Helvetica", "", size)
    pdf.multi_cell(0, 5, sanitize(t), **MC); pdf.set_text_color(*INK)


pdf = PDF(format="A4"); pdf.set_auto_page_break(True, 16); pdf.add_page()
pdf.ln(6)
pdf.set_text_color(*INK); pdf.set_font("Helvetica", "B", 20)
pdf.multi_cell(0, 8, sanitize(TAGLINE), **MC)
pdf.ln(1); body(pdf, SUB, 10.5)
pdf.ln(2)
# ROAI callout
y = pdf.get_y(); pdf.set_fill_color(*LITE); pdf.rect(14, y, 182, 16, "F")
pdf.set_xy(18, y + 3); pdf.set_text_color(*INK); pdf.set_font("Courier", "B", 11)
pdf.cell(0, 5, sanitize(ROAI_F), 0, 1)
pdf.set_xy(18, y + 9); pdf.set_text_color(*GREY); pdf.set_font("Helvetica", "", 8.5)
pdf.cell(0, 5, sanitize(ROAI_NOTE), 0, 1)
pdf.set_y(y + 18)

h2(pdf, "The opportunity")
for n, l in STATS:
    pdf.set_font("Helvetica", "B", 11); pdf.set_text_color(*IND); pdf.cell(20, 5, sanitize(n), 0, 0)
    pdf.set_font("Helvetica", "", 10); pdf.set_text_color(*GREY); pdf.cell(0, 5, sanitize(l), 0, 1)
pdf.set_text_color(*INK)

h2(pdf, "One platform, two modes")
pdf.set_font("Helvetica", "B", 10); pdf.cell(0, 5, "Agency command center", 0, 1)
body(pdf, "Portfolio ROAI, revenue & margin - client directory & health tiers - ROAI analytics - billing with live Metronome usage data.")
pdf.set_font("Helvetica", "B", 10); pdf.set_text_color(*INK); pdf.cell(0, 5, "Customer portal", 0, 1)
body(pdf, "A personal ROAI Center - active projects across six pillars - AI marketplace to launch agents - self-serve billing, usage & top-ups.")

h2(pdf, "Six CX pillars, powered by the ON.Ecosystem")
for name, desc in PILLARS:
    pdf.set_font("Helvetica", "B", 9.5); pdf.set_text_color(*INK); pdf.cell(48, 5, sanitize(name), 0, 0)
    pdf.set_font("Helvetica", "", 9.5); pdf.set_text_color(*GREY); pdf.cell(0, 5, sanitize(desc), 0, 1)
pdf.set_font("Helvetica", "I", 9); pdf.set_text_color(*IND)
pdf.multi_cell(0, 5, sanitize("Agents: " + "  -  ".join(AGENTS)), **MC)
pdf.set_text_color(*INK)

h2(pdf, "Live usage-based billing (Metronome)")
body(pdf, "The Live panel reads a real Metronome account via a secure read-only proxy (the key never reaches the browser): 7 customers on prepaid-commit or hybrid models, finalized + draft invoices with 6 months of revenue history, commit/credit drawdown, and a full catalog of metrics, products, rate cards and packages.")

# page 2: proof + demo
pdf.add_page()
h2(pdf, "Proof - real FPT outcomes behind the demo")
for sol, head, who, metrics in CASES:
    pdf.set_font("Helvetica", "B", 11); pdf.set_text_color(*INK); pdf.cell(0, 6, sanitize(head), 0, 1)
    pdf.set_font("Helvetica", "", 9); pdf.set_text_color(*IND); pdf.cell(0, 4.5, sanitize(sol + "   |   " + who), 0, 1)
    pdf.set_text_color(*GREY); pdf.set_font("Helvetica", "", 9); pdf.multi_cell(0, 4.5, sanitize(metrics), **MC); pdf.ln(1)
pdf.set_font("Helvetica", "", 9.5); pdf.set_text_color(*INK)
pdf.multi_cell(0, 5, sanitize("Partners: " + "  -  ".join(PARTNERS)), **MC)
pdf.set_font("Helvetica", "B", 10); pdf.multi_cell(0, 5, "1,000+ certified engineers      1,500+ platform certifications", **MC)

h2(pdf, "Run the demo - 7 steps (~6-8 min)")
for i, (head, d, u) in enumerate(DEMO):
    pdf.set_font("Helvetica", "B", 9.5); pdf.set_text_color(*IND); pdf.cell(7, 5, str(i + 1) + ".", 0, 0)
    pdf.set_text_color(*INK); pdf.cell(50, 5, sanitize(head), 0, 0)
    pdf.set_font("Helvetica", "", 9); pdf.set_text_color(*GREY); pdf.multi_cell(0, 5, sanitize(d), **MC)
pdf.ln(2)
y = pdf.get_y(); pdf.set_fill_color(*LITE); pdf.rect(14, y, 182, 18, "F")
pdf.set_xy(18, y + 3.5); pdf.set_text_color(*INK); pdf.set_font("Helvetica", "B", 10)
pdf.cell(0, 5, "See it live", 0, 1)
pdf.set_xy(18, y + 9); pdf.set_font("Helvetica", "", 9.5); pdf.set_text_color(*IND)
pdf.cell(0, 5, sanitize("Demo:  " + LIVE.replace("https://", "") + "      Guide:  " + GUIDE.replace("https://", "")), 0, 1)

pdf_path = os.path.join(OUT, "Vela-OS-One-Pager.pdf")
pdf.output(pdf_path)
print("wrote", pdf_path)
